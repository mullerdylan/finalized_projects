from pathlib import WindowsPath
from sys import setswitchinterval
from tkinter import *
import datetime
import os.path
import tkinter
from xlsxwriter.workbook import Workbook
from tkinter import ttk
from win32com.client import Dispatch
from tkinter.filedialog import FileDialog, askdirectory, askopenfilename
from PIL import ImageTk, Image 
from pptx import Presentation
import xlrd
from pptx.util import Inches
import pyqrcode 
import png 
from pyqrcode import QRCode
import os 
from pptx.util import Inches, Pt
import tkinter.messagebox

root = Tk()
root.geometry("960x540")
root.title("USPS Sign Creator")
root.iconbitmap('usps-icon.ico') 
bg = ImageTk.PhotoImage(file = "uspslogo.png")
my_canvas = Canvas(root, width = 960, height =540, bd = 0, highlightthickness=0)
my_canvas.pack(fill = "both", expand = True)

#set image in canvas 
my_canvas.create_image(0,0, image = bg, anchor = "nw")

#Creates static text
my_canvas.create_text(500,165, text= "enter in value in AGT00G018025046PIC format ", anchor = 'c',font = ("helvetica",15) ,fill= "white")
my_canvas.create_text(10,520, text= "Tell Dylan he is bad at his job:  dylan.muller@usps.gov", anchor = "nw", font = ("helvetica",8) ,fill= "white")
vartext  = my_canvas.create_text(500,60, text= "", anchor = "c", font = ("helvetica",25) ,fill= "white")
EINLvartext  = my_canvas.create_text(25,351, text= "", anchor = "nw", font = ("helvetica",15) ,fill= "white")
Durvartext  = my_canvas.create_text(25,391, text= "", anchor = "nw", font = ("helvetica",15) ,fill= "white")

#________________set file defaults
impGP8 = ("GP82.PNG")
impGT10 = ("GT10.PNG")
header = ("Header.png")

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]

#Creates input box
IB = Entry(root, font =("Helvetica", 24), width = 20, fg = "black", bd = 0)
#adds boxes to canvas
IBwindow = my_canvas.create_window(500,120,window= IB)

#sets focus for input
IB.focus_set()

def entryget(event):
    global val, strt
    val = IB.get().upper()
    if val[0:1] == 'A' and len(val) == 18 and val[5:6].isalpha():
        signcreate()
        IB.delete(0,END)
        strt = datetime.datetime.now()
    else:
        tkinter.messagebox.showerror(title='Invalid Input: Check String Length', message='Input text in AGT00G018025046PIC format')

def bulkimpport():
    global val
    file_path = askopenfilename()
    extension = os.path.splitext(file_path)[1] 
    filetemp = os.path.basename(file_path)
    fin_filename = os.path.splitext(filetemp)[0]

    if extension == '.xls':
        wb = xlrd.open_workbook(file_path)
        sheet = wb.sheet_by_index(0)
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        strt = datetime.datetime.now()
        for i  in range(1,sheet.nrows):
        #Generates QR Code     
            #assignes variable of BC to barcode string
            BC = sheet.cell_value(i , 1) 
            #if row 3 indicates GT10, assign AGT to prefix     
            if sheet.cell_value(i, 2) == 'GT10':
                val = "AGT" + BC + "PIC"
                usepic = impGT10
            #if row 3 indicates GP8, assign AGP to prefix 
            else:
                val = "AGP" + BC + "PIC"
                usepic = impGP8
            url = pyqrcode.create(val)
            url.png('%s.png' % val, scale = 13, quiet_zone=2)
            #_________starts sign creatioin loop
            BCR =  val
            url = pyqrcode.create(BCR)
            url.png('%s.png' % BCR, scale = 13, quiet_zone=2)
                #END QR CODE GENERATION 

            #SLIDE TEMPLATE CREATION 
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            txBox = slide.shapes.add_textbox(Inches(4.5), Inches(6), Inches(1), Inches(1))

            tf = txBox.text_frame
            p = tf.add_paragraph()

            #locates position of imported images
            
            slide.shapes.add_picture(usepic, Inches(.25), Inches(1))  
            slide.shapes.add_picture(header, Inches(2.15),0)  
            slide.shapes.add_picture('%s.png' % BCR, Inches(4.5), Inches(1)) 

            #ADDING dynamic text
            #left = top = width = height 
            p.text = "%s" % BCR
            p.font.size = Pt(40)
            p.font.bold = False

            #Deletes QR CODE
            os.remove("%s.png" % BCR) 
        #saves file and records timestamp
        prs.save('Signs_' + '%s.pptx' %fin_filename)
        nd = datetime.datetime.now()
        dur = nd - strt
        my_canvas.itemconfig(Durvartext, text= 'Signs Created, file Name is: {}.pptx Duration: {}'.format(fin_filename,dur), fill= 'light green' )
    else: 
        tkinter.messagebox.showerror(title='Invalid Input', message='File needs to be in .xls format')


def signcreate():
    strt = datetime.datetime.now()
    BCR =  val
    url = pyqrcode.create(BCR)
    url.png('%s.png' % BCR, scale = 13, quiet_zone=2)
    #END QR CODE GENERATION 

    #SLIDE TEMPLATE CREATION 
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(Inches(4.5), Inches(6), Inches(1), Inches(1))

    tf = txBox.text_frame
    p = tf.add_paragraph()

    #locates position of imported images
    btype = val[0:3]
    if btype == 'AGT':
        slide.shapes.add_picture(impGT10, Inches(.25), Inches(1))  
    else:
        slide.shapes.add_picture(impGP8, Inches(.25), Inches(1)) 
    slide.shapes.add_picture(header, Inches(2.15),0)  
    slide.shapes.add_picture('%s.png' % BCR, Inches(4.5), Inches(1)) 

    #ADDING dynamic text
    #left = top = width = height 
    p.text = "%s" % BCR
    p.font.size = Pt(40)
    p.font.bold = False

    #Deletes QR CODE
    os.remove("%s.png" % BCR) 
    prs.save('Signs_' + '%s.pptx' %val)
    nd = datetime.datetime.now()
    dur = nd - strt
    my_canvas.itemconfig(Durvartext, text= 'Sign Created for {} Duration: {}'.format(val,dur), fill= 'light green' )
    
Bulksign_bbtn = Button(root, text = "Bulk Sign Creation (.xls)", width=25, anchor='center', command = bulkimpport)
query_window = my_canvas.create_window(850,520,window= Bulksign_bbtn)

root.bind('<Return>', entryget)
root.mainloop()
